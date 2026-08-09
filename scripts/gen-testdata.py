#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = ["reportlab", "python-docx", "openpyxl", "python-pptx", "pillow>=10.1"]
# ///
"""Generate the binary test fixtures in testdata/.

The fixtures are committed, so this script exists to make them reproducible and
to document what each one is *for*. Run it only when a fixture needs to change:

    ./scripts/gen-testdata.py

Text fixtures (.md, .txt, .csv) are hand-written and committed directly; this
script does not touch them.

The PDFs are byte-identical run to run -- `invariant=1` stops reportlab writing
a timestamp and a random document ID, so regenerating them produces no diff and
a real change is visible. The OOXML fixtures are not: python-docx, openpyxl and
python-pptx stamp times into the package, so sample.docx, sample.xlsx and
sample.pptx differ on every run even when nothing about them changed. Check
their diffs by hand or discard them.

The PDFs are the interesting ones, because they are what exercise routing:

  text.pdf           every page has real text operators  -> native extraction
  scanned.pdf        image only, no text operators       -> must route to OCR
  mixed.pdf          one of each                         -> must route per page
  scanned-table.pdf  a ruled table drawn as pixels       -> needs layout analysis
  faded.pdf          a scan too faint for OCR            -> must escalate to vision
  corrupt.pdf        a PDF header over garbage           -> must fail cleanly

scanned-table.pdf is the one that separates the OCR tiers: text-line OCR reads
the cells but returns them as flat paragraphs, and only layout analysis
recovers the grid.

faded.pdf separates the OCR tiers from the vision tier, and its parameters are
measured rather than chosen -- see `degrade()`.
"""

import io
import json
import pathlib
import random

from docx import Document as Docx
from docx.shared import Pt
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

OUT = pathlib.Path(__file__).resolve().parent.parent / "testdata"

# ---------------------------------------------------------------------------
# Fixture content
#
# Declared once and used twice: to draw the fixture, and to write the ground
# truth the benchmark scores against. Because these documents are generated
# rather than collected, the expected output is known exactly rather than
# transcribed by hand -- which is what makes a character error rate meaningful
# on them at all.
# ---------------------------------------------------------------------------

TEXT_PDF_PAGES = [
    (
        "Quarterly Report",
        [
            "Revenue grew by twelve percent over the prior quarter.",
            "Operating costs were flat.",
            "The outlook for the coming quarter remains positive.",
        ],
    ),
    (
        "Appendix A",
        [
            "All figures are unaudited.",
            "Currency is United States dollars.",
        ],
    ),
]

SCANNED_LINES = [
    "SCANNED INVOICE",
    "",
    "Invoice number 4471",
    "Amount due 1,250.00",
    "Due on receipt",
]

MIXED_TEXT_PAGE = (
    "Cover Letter",
    [
        "Please find the signed agreement attached.",
        "The scanned copy follows on the next page.",
    ],
)

MIXED_SCANNED_LINES = [
    "SIGNED AGREEMENT",
    "",
    "Party A and Party B agree",
    "Signature on file",
]

def flatten(rows) -> list[str]:
    """Table cells in reading order, for the text expectation."""
    return [cell for row in rows for cell in row]


# sample.csv is hand-written and committed; these are its contents, restated so
# the ground truth can be generated from one place.
CSV_ROWS = [
    ("region", "units", "revenue"),
    ("North", "120", "14400.00"),
    ("South", "86", "10320.00"),
    ("East", "203", "24360.00"),
    ("West", "54", "6480.00"),
]

DOCX_ROUTING_TABLE = [
    ("Format", "Engine", "Needs OCR"),
    ("DOCX", "anydoc", "no"),
    ("Scanned PDF", "paddleocr", "yes"),
]

# A merged cell: the second column of row 0 is a shadow slot, which the
# benchmark skips, so that row has one cell where the other has two.
DOCX_MERGED_TABLE = [
    ("Spans two columns",),
    ("left", "right"),
]

# A faded page: the content is ordinary, and everything interesting about the
# fixture is in how it is rendered.
FADED_LINES = [
    "SHIPPING RECEIPT",
    "",
    "Consignment 8842-QX",
    "Received 14 March 2026",
    "",
    "Freight charge 1,284.50",
    "Fuel surcharge 96.25",
    "Handling 40.00",
    "",
    "Total due 1,420.75",
    "Terms net thirty days",
]

TABLE_TITLE = "QUARTERLY SALES"
TABLE_NOTE = "All figures are unaudited."
TABLE_FOOTER = "Totals exclude tax."
TABLE_ROWS = [
    ("Region", "Units", "Revenue"),
    ("North", "120", "14,400.00"),
    ("South", "86", "10,320.00"),
    ("East", "203", "24,360.00"),
    ("West", "54", "6,480.00"),
]


def page_image(text: str, size=(1200, 1600)) -> Image.Image:
    """A white page with black text, drawn as pixels.

    This is what makes scanned.pdf a genuine OCR case: the words are visible to
    a human and to an OCR engine, and completely absent from the PDF's content
    stream.
    """
    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    y = 120
    for line in text.splitlines():
        draw.text((120, y), line, fill="black")
        y += 60
    return img


def text_page(c: canvas.Canvas, title: str, body: list[str]) -> None:
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, 720, title)
    c.setFont("Helvetica", 11)
    y = 690
    for line in body:
        c.drawString(72, y, line)
        y -= 16


def image_page(c: canvas.Canvas, img: Image.Image, **opts) -> None:
    buf = io.BytesIO()
    img.save(buf, format="PNG", **opts)
    buf.seek(0)
    c.drawImage(ImageReader(buf), 0, 0, width=LETTER[0], height=LETTER[1])


def write_text_pdf() -> None:
    c = canvas.Canvas(str(OUT / "text.pdf"), pagesize=LETTER, invariant=1)
    for title, body in TEXT_PDF_PAGES:
        text_page(c, title, body)
        c.showPage()
    c.save()


def write_scanned_pdf() -> None:
    c = canvas.Canvas(str(OUT / "scanned.pdf"), pagesize=LETTER, invariant=1)
    image_page(c, page_image("\n".join(SCANNED_LINES)))
    c.showPage()
    c.save()


def write_mixed_pdf() -> None:
    c = canvas.Canvas(str(OUT / "mixed.pdf"), pagesize=LETTER, invariant=1)
    text_page(c, MIXED_TEXT_PAGE[0], MIXED_TEXT_PAGE[1])
    c.showPage()
    image_page(c, page_image("\n".join(MIXED_SCANNED_LINES)))
    c.showPage()
    c.save()


def table_image(size=(1700, 2200)) -> Image.Image:
    """A ruled table drawn as pixels, with no text operators anywhere.

    Deliberately *wired* -- every cell fully bordered -- because that is the
    case a table-structure model should get right, and because a human reading
    the rendered page can see immediately whether the recovered grid matches.
    """
    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)

    rows = TABLE_ROWS
    left, top = 200, 400
    col_w, row_h = 380, 120

    draw.text((left, top - 160), TABLE_TITLE, fill="black")
    draw.text((left, top - 100), TABLE_NOTE, fill="black")

    for r in range(len(rows) + 1):
        y = top + r * row_h
        draw.line([(left, y), (left + col_w * 3, y)], fill="black", width=3)
    for c in range(4):
        x = left + c * col_w
        draw.line([(x, top), (x, top + row_h * len(rows))], fill="black", width=3)

    for r, cells in enumerate(rows):
        for c, value in enumerate(cells):
            draw.text((left + c * col_w + 30, top + r * row_h + 45), value, fill="black")

    draw.text((left, top + row_h * len(rows) + 80), TABLE_FOOTER, fill="black")
    return img


def write_scanned_table_pdf() -> None:
    c = canvas.Canvas(str(OUT / "scanned-table.pdf"), pagesize=LETTER, invariant=1)
    image_page(c, table_image())
    c.showPage()
    c.save()


# The degradation parameters below are measured, not chosen. They came from a
# sweep against the real engines, and they sit in a narrow band: one step
# gentler and PaddleOCR reads the page confidently and wrongly, one step harsher
# and MinerU loses it too. Measured on this corpus, one page:
#
#   parameters      OCR chars  OCR conf   OCR CER   vision CER
#   (clean)               157     0.988     0.019          --
#   one step gentler      152     0.918     0.100       0.006   <- no escalation
#   THESE                   1     0.495     1.000       0.019   <- escalates
#   one step harsher        0        --     1.000       1.000   <- both fail
#
# The gentler row is why this fixture is set where it is, and it is also the
# honest limit of per-page quality scoring: OCR misread a tenth of that page
# and reported 92% confidence, so nothing downstream can tell.
FADED = dict(
    scale=0.8,     # rendered small and blown back up: strokes lose their edges
    blur=0.9,
    ink=170,       # near-white paper, barely-grey ink: an exhausted photocopier
    paper=190,
    noise=12,
    speckle=0.004,
    angle=-1.5,    # fed in slightly crooked
    quality=28,    # and scanned to a low-quality JPEG
)

# Fixed so regeneration is reproducible. PIL's own effect_noise() takes no seed.
FADED_SEED = 20260808


def faded_image(size=(1700, 2200), font_size=46) -> Image.Image:
    """The faded fixture: legible to a person, unreadable to OCR.

    That combination is the whole point. A page OCR cannot read because there
    is nothing on it proves nothing; this one carries its text plainly enough
    that a human reads it at a glance and the vision tier recovers it almost
    exactly, while the OCR tier returns a single character.
    """
    img = Image.new("L", size, 255)
    draw = ImageDraw.Draw(img)
    # load_default(size=...) rather than the bitmap default used elsewhere:
    # this fixture needs text large enough that failing to read it is about
    # contrast rather than about resolution.
    font = ImageFont.load_default(size=font_size)
    y = 300
    for line in FADED_LINES:
        if line:
            draw.text((220, y), line, fill=0, font=font)
        y += int(font_size * 1.9)
    return degrade(img, **FADED)


def degrade(img, *, scale, blur, ink, paper, noise, speckle, angle, quality):
    """Age a rendered page the way a bad scan does."""
    rng = random.Random(FADED_SEED)
    w, h = img.size

    if scale < 1:
        img = img.resize((int(w * scale), int(h * scale)), Image.BILINEAR)
        img = img.resize((w, h), Image.BILINEAR)
    img = img.filter(ImageFilter.GaussianBlur(blur))
    # Compress the dynamic range: black moves up toward grey, white moves down.
    img = img.point(lambda v: int(ink + (paper - ink) * (v / 255)))
    grain = Image.frombytes(
        "L", img.size,
        bytes(max(0, min(255, int(rng.gauss(128, noise)))) for _ in range(w * h)),
    )
    img = Image.blend(img, grain, 0.35)
    px = img.load()
    for _ in range(int(w * h * speckle)):
        px[rng.randrange(w), rng.randrange(h)] = rng.choice((0, 255))
    img = img.rotate(angle, resample=Image.BILINEAR, fillcolor=paper)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality)
    out = Image.open(io.BytesIO(buf.getvalue()))
    out.load()
    # Grayscale, not RGB: the page has no colour, and three identical channels
    # would triple what the PNG has to store.
    return out.convert("L")


def write_faded_pdf() -> None:
    # PNG, and grayscale, and at this exact size -- all three measured rather
    # than preferred, because this fixture is fragile in a way the others are
    # not. Its whole content is a 20-level contrast range, so:
    #
    #   - a JPEG on the way into the PDF, even at quality 95, quantizes that
    #     range away and the vision tier stops reading the page too;
    #   - rendering smaller (1275x1650, 1020x1320) loses both tiers outright;
    #   - posterizing to compress the noise moves both tiers unpredictably --
    #     at 4 bits OCR recovers to CER 0.36, which is not this fixture.
    #
    # So it costs about 900KB, which is what a page of noise costs when it has
    # to stay exactly as noisy as it was measured to be.
    c = canvas.Canvas(str(OUT / "faded.pdf"), pagesize=LETTER, invariant=1)
    image_page(c, faded_image(), optimize=True)
    c.showPage()
    c.save()


def write_corrupt_pdf() -> None:
    # A valid header over bytes that are not a PDF body: the shim must report
    # this as malformed rather than panicking or returning an empty document.
    (OUT / "corrupt.pdf").write_bytes(b"%PDF-1.7\n" + b"\x00\xff not a pdf body \x01" * 40)


def write_docx() -> None:
    d = Docx()
    d.add_heading("Engineering Handbook", level=1)
    d.add_paragraph("This document exercises headings, lists, tables and styling.")

    d.add_heading("Principles", level=2)
    for item in ("Extract natively whenever possible.", "OCR only pages that require OCR."):
        d.add_paragraph(item, style="List Bullet")
    for item in ("Inspect the document.", "Route each page.", "Normalize the result."):
        d.add_paragraph(item, style="List Number")

    d.add_heading("Styling", level=2)
    p = d.add_paragraph("Text can be ")
    p.add_run("bold").bold = True
    p.add_run(", ")
    p.add_run("italic").italic = True
    p.add_run(", or plain.")

    d.add_heading("Routing table", level=2)
    t = d.add_table(rows=3, cols=3)
    t.style = "Table Grid"
    for row, cells in enumerate(DOCX_ROUTING_TABLE):
        for col, value in enumerate(cells):
            cell = t.cell(row, col)
            cell.text = value
            if row == 0:
                cell.paragraphs[0].runs[0].font.bold = True
                cell.paragraphs[0].runs[0].font.size = Pt(11)

    # A merged cell, so the canonical grid's covered_by slots get exercised.
    t2 = d.add_table(rows=2, cols=2)
    t2.style = "Table Grid"
    t2.cell(0, 0).merge(t2.cell(0, 1)).text = "Spans two columns"
    t2.cell(1, 0).text = "left"
    t2.cell(1, 1).text = "right"

    d.save(OUT / "sample.docx")


def write_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Inventory"
    for row in [("item", "qty", "price"), ("widget", 3, 9.99), ("gadget", 12, 24.50)]:
        ws.append(row)

    ws2 = wb.create_sheet("Notes")
    ws2.append(("note",))
    ws2.append(("Prices exclude tax.",))
    wb.save(OUT / "sample.xlsx")


def write_pptx() -> None:
    prs = Presentation()
    title_layout, bullet_layout = prs.slide_layouts[0], prs.slide_layouts[1]

    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Document Processing"
    s1.placeholders[1].text = "A per-page routing pipeline"

    s2 = prs.slides.add_slide(bullet_layout)
    s2.shapes.title.text = "Pipeline"
    body = s2.placeholders[1].text_frame
    body.text = "Inspect"
    for step in ("Route", "Extract", "Normalize"):
        body.add_paragraph().text = step

    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Embedded image"
    buf = io.BytesIO()
    page_image("chart placeholder", size=(600, 400)).save(buf, format="PNG")
    buf.seek(0)
    s3.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(4))

    prs.save(OUT / "sample.pptx")


def write_ground_truth() -> None:
    """Write what each fixture is supposed to say.

    `scripts/bench.py` scores extraction against this. Because these documents
    are generated rather than collected, the expectation is exact rather than
    transcribed, so a character error rate computed against it means something.

    What it does *not* mean is real-world accuracy: these are clean synthetic
    renderings, not photographs of creased paper. The format is deliberately
    corpus-agnostic so a directory of real documents with hand-written ground
    truth can be scored by the same harness.
    """
    truth = {
        "text.pdf": {
            "pages": [
                {"number": n, "text": [title, *body]}
                for n, (title, body) in enumerate(TEXT_PDF_PAGES, start=1)
            ]
        },
        "scanned.pdf": {
            "pages": [{"number": 1, "text": [ln for ln in SCANNED_LINES if ln]}]
        },
        # The OCR tiers score close to 1.0 CER on this one. That is the
        # expected result, not a broken expectation: it is what the vision
        # tier's benchmark row is measured against.
        "faded.pdf": {
            "pages": [{"number": 1, "text": [ln for ln in FADED_LINES if ln]}]
        },
        "mixed.pdf": {
            "pages": [
                {"number": 1, "text": [MIXED_TEXT_PAGE[0], *MIXED_TEXT_PAGE[1]]},
                {"number": 2, "text": [ln for ln in MIXED_SCANNED_LINES if ln]},
            ]
        },
        "scanned-table.pdf": {
            "pages": [
                {
                    # In reading order, table contents included: the text score
                    # measures whether the characters were read, and an
                    # engine's inability to structure them should not also
                    # register as a failure to read them.
                    "number": 1,
                    "text": [TABLE_TITLE, TABLE_NOTE, *flatten(TABLE_ROWS), TABLE_FOOTER],
                    # Scored separately, because recovering a grid is a
                    # different capability that only the layout tier claims.
                    "tables": [[list(row) for row in TABLE_ROWS]],
                }
            ]
        },
        "sample.csv": {
            "pages": [
                {
                    "number": 1,
                    "text": flatten(CSV_ROWS),
                    "tables": [[list(row) for row in CSV_ROWS]],
                }
            ]
        },
        "sample.docx": {
            "pages": [
                {
                    "number": 1,
                    "text": [
                        "Engineering Handbook",
                        "This document exercises headings, lists, tables and styling.",
                        "Principles",
                        "Extract natively whenever possible.",
                        "OCR only pages that require OCR.",
                        "Inspect the document.",
                        "Route each page.",
                        "Normalize the result.",
                        "Styling",
                        "Text can be bold, italic, or plain.",
                        "Routing table",
                        *flatten(DOCX_ROUTING_TABLE),
                        *flatten(DOCX_MERGED_TABLE),
                    ],
                }
            ]
        },
        "sample.txt": {
            "pages": [
                {
                    "number": 1,
                    "text": [
                        "Plain text carries no markup, and nothing here should be interpreted as any.",
                        "Characters like # and * and | are literal in this file. A line that looks",
                        "like - this is not a list item.",
                        "Indented lines keep their indentation.",
                        "The last paragraph has no trailing newline issues.",
                    ],
                }
            ]
        },
    }
    # The DOCX tables are expected too: the native path recovers grids as much
    # as the layout tier does, and a regression there should show up here.
    truth["sample.docx"]["pages"][0]["tables"] = [
        [list(row) for row in DOCX_ROUTING_TABLE],
        [list(row) for row in DOCX_MERGED_TABLE],
    ]
    path = OUT / "ground-truth.json"
    path.write_text(json.dumps(truth, indent=2, ensure_ascii=False) + "\n")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    write_text_pdf()
    write_scanned_pdf()
    write_scanned_table_pdf()
    write_mixed_pdf()
    write_faded_pdf()
    write_corrupt_pdf()
    write_docx()
    write_xlsx()
    write_pptx()
    write_ground_truth()
    for path in sorted(OUT.glob("*")):
        if path.is_file():
            print(f"{path.name:16} {path.stat().st_size:>8} bytes")


if __name__ == "__main__":
    main()
